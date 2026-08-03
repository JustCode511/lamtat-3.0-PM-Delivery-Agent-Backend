"""
Talent Report Generator — builds PPTX / XLSX / DOCX reports from talent data.

Uses python-pptx, openpyxl, and python-docx (all pure Python — no Office needed).

Report structure (mirrors the PM PPT layout):
  1. Title slide / header
  2. Team overview (counts by status)
  3. Availability breakdown (per-employee flag)
  4. Skill coverage (top skills with holder counts)
  5. Active projects (name, client, team)
  6. Rolling off in next 30 days
"""
from __future__ import annotations

import io
from collections import Counter
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter

from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ---------------------------------------------------------------------------
# Brand palette (matches PM report)
# ---------------------------------------------------------------------------
C_DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
C_MED_BLUE  = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT_BG  = RGBColor(0xD6, 0xE4, 0xF0)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN     = RGBColor(0x70, 0xAD, 0x47)
C_AMBER     = RGBColor(0xFF, 0xC0, 0x00)
C_RED       = RGBColor(0xC0, 0x00, 0x00)
C_DARK_TEXT = RGBColor(0x26, 0x26, 0x26)


# ---------------------------------------------------------------------------
# Shared context builder — derives every metric each generator needs
# ---------------------------------------------------------------------------

def _build_context(employees, projects, allocations) -> dict[str, Any]:
    """Compute the summary structures the three generators share."""
    today = date.today()

    total = len(employees)
    counts = Counter(e.status for e in employees)
    available = counts.get("available", 0)
    allocated = counts.get("allocated", 0)
    bench = counts.get("bench", 0)
    on_leave = counts.get("on_leave", 0)

    # Rolling off in next 30 days
    rolling_off: list[dict[str, Any]] = []
    for e in employees:
        if e.status != "allocated":
            continue
        try:
            d = date.fromisoformat(e.availability_date)
            days = (d - today).days
        except (ValueError, TypeError):
            continue
        if 0 <= days <= 30:
            rolling_off.append({
                "name": e.name,
                "designation": e.designation,
                "current_project": e.current_project or "-",
                "days_until_free": days,
                "availability_date": e.availability_date,
            })
    rolling_off.sort(key=lambda r: r["days_until_free"])

    # Skill coverage — top 15 primary skills
    skill_counter: Counter = Counter()
    for e in employees:
        for s in e.primary_skills:
            skill_counter[s] += 1
    top_skills = skill_counter.most_common(15)

    active_projects = [p for p in projects if p.status in ("active", "planning")]

    return {
        "today": today,
        "total_employees": total,
        "available": available,
        "allocated": allocated,
        "bench": bench,
        "on_leave": on_leave,
        "rolling_off": rolling_off,
        "top_skills": top_skills,
        "active_projects": active_projects,
        "employees": employees,
        "projects": projects,
        "allocations": allocations,
    }


# ═══════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════

def _pptx_textbox(slide, left, top, width, height, text, *,
                  font_size=12, bold=False, color=C_DARK_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _pptx_slide_header(slide, title: str) -> None:
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = C_DARK_BLUE
    header.line.fill.background()
    tf = header.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = C_WHITE


def _pptx_table(slide, left, top, width, height, headers, rows) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BLUE
        tf = cell.text_frame
        tf.text = h
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = C_WHITE

    # body rows (alternating bg)
    for i, row in enumerate(rows, start=1):
        alt = i % 2 == 0
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            if alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_BG
            tf = cell.text_frame
            tf.text = str(val)
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = C_DARK_TEXT


def generate_talent_pptx_bytes(employees, projects, allocations) -> bytes:
    ctx = _build_context(employees, projects, allocations)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BLUE
    bg.line.fill.background()
    _pptx_textbox(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
                  "Talent Management Report", font_size=44, bold=True,
                  color=C_WHITE, align=PP_ALIGN.CENTER)
    _pptx_textbox(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.6),
                  f"Team Overview · {ctx['today'].isoformat()}", font_size=18,
                  color=C_LIGHT_BG, align=PP_ALIGN.CENTER)

    # ── Slide 2: Team Overview ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _pptx_slide_header(s, "Team Overview")

    metrics = [
        ("Total Employees", ctx["total_employees"], C_MED_BLUE),
        ("Available", ctx["available"], C_GREEN),
        ("Allocated", ctx["allocated"], C_MED_BLUE),
        ("Bench", ctx["bench"], C_AMBER),
        ("On Leave", ctx["on_leave"], C_RED),
    ]
    box_w = Inches(2.4)
    gap = Inches(0.15)
    total_w = box_w * len(metrics) + gap * (len(metrics) - 1)
    start_left = (prs.slide_width - total_w) / 2
    top = Inches(1.4)
    for i, (label, value, color) in enumerate(metrics):
        left = start_left + (box_w + gap) * i
        box = s.shapes.add_shape(1, left, top, box_w, Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        _pptx_textbox(s, left, top + Inches(0.15), box_w, Inches(0.7),
                      str(value), font_size=40, bold=True, color=C_WHITE,
                      align=PP_ALIGN.CENTER)
        _pptx_textbox(s, left, top + Inches(0.95), box_w, Inches(0.4),
                      label, font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 3: Availability breakdown (first 15) ──────────────────
    s = prs.slides.add_slide(blank)
    _pptx_slide_header(s, "Availability Breakdown")
    rows = []
    for e in ctx["employees"][:15]:
        flag = _availability_flag(e, ctx["today"])
        rows.append((e.name, e.designation, e.status.upper(), flag, e.availability_date))
    _pptx_table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.7),
                ["Name", "Designation", "Status", "Availability", "Date"], rows)

    # ── Slide 4: Skill Coverage ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _pptx_slide_header(s, "Skill Coverage")
    rows = [(name, count) for name, count in ctx["top_skills"]]
    _pptx_table(s, Inches(3.0), Inches(1.1), Inches(7.3), Inches(5.7),
                ["Skill", "Primary Holders"], rows)

    # ── Slide 5: Active Projects + Rolling Off ──────────────────────
    s = prs.slides.add_slide(blank)
    _pptx_slide_header(s, "Active Projects & Upcoming Availability")

    _pptx_textbox(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.4),
                  "Active Projects", font_size=15, bold=True, color=C_DARK_BLUE)
    proj_rows = [
        (p.name, p.client, p.status.upper(), p.end_date)
        for p in ctx["active_projects"][:8]
    ]
    _pptx_table(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3),
                ["Project", "Client", "Status", "Ends"], proj_rows or [("-", "-", "-", "-")])

    _pptx_textbox(s, Inches(7.0), Inches(1.05), Inches(5.8), Inches(0.4),
                  "Rolling Off in 30 Days", font_size=15, bold=True, color=C_DARK_BLUE)
    ro_rows = [
        (r["name"], r["current_project"], f"{r['days_until_free']}d")
        for r in ctx["rolling_off"][:8]
    ]
    _pptx_table(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
                ["Name", "Project", "In"], ro_rows or [("None", "-", "-")])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _availability_flag(e, today: date) -> str:
    try:
        d = date.fromisoformat(e.availability_date)
        days_until = (d - today).days
    except (ValueError, TypeError):
        days_until = 9999
    if e.status == "available":
        return "AVAILABLE NOW"
    if e.status == "bench":
        return "BENCH (available now)"
    if e.status == "on_leave":
        return f"ON LEAVE until {e.availability_date}"
    if e.status == "allocated" and days_until <= 30:
        return f"FREE IN {days_until}d"
    return f"ALLOCATED ({e.allocation_pct}%)"


# ═══════════════════════════════════════════════════════════════════════════
# XLSX
# ═══════════════════════════════════════════════════════════════════════════

_HEADER_FILL = PatternFill("solid", fgColor="1F3864")
_HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
_ZEBRA_FILL = PatternFill("solid", fgColor="D6E4F0")


def _xlsx_write_header(ws, headers: list[str]) -> None:
    ws.append(headers)
    for j in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=j)
        cell.fill = _HEADER_FILL
        cell.font = _HEADER_FONT
        cell.alignment = Alignment(horizontal="left", vertical="center")


def _xlsx_autosize(ws) -> None:
    for col in ws.columns:
        max_len = max((len(str(c.value)) if c.value is not None else 0) for c in col)
        ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 2, 40)


def generate_talent_xlsx_bytes(employees, projects, allocations) -> bytes:
    ctx = _build_context(employees, projects, allocations)
    wb = Workbook()

    # -- Sheet 1: Overview --
    ws = wb.active
    ws.title = "Overview"
    ws["A1"] = "Talent Management Report"
    ws["A1"].font = Font(bold=True, size=16, color="1F3864")
    ws["A2"] = f"Generated on {ctx['today'].isoformat()}"
    ws["A2"].font = Font(italic=True, size=10)

    ws["A4"] = "Metric"
    ws["B4"] = "Value"
    for j in (1, 2):
        c = ws.cell(row=4, column=j)
        c.fill = _HEADER_FILL
        c.font = _HEADER_FONT
    metrics = [
        ("Total Employees", ctx["total_employees"]),
        ("Available", ctx["available"]),
        ("Allocated", ctx["allocated"]),
        ("Bench", ctx["bench"]),
        ("On Leave", ctx["on_leave"]),
        ("Rolling Off in 30d", len(ctx["rolling_off"])),
        ("Active Projects", len(ctx["active_projects"])),
    ]
    for i, (k, v) in enumerate(metrics, start=5):
        ws.cell(row=i, column=1, value=k)
        ws.cell(row=i, column=2, value=v)
    ws.column_dimensions["A"].width = 24
    ws.column_dimensions["B"].width = 14

    # -- Sheet 2: Employees --
    ws2 = wb.create_sheet("Employees")
    _xlsx_write_header(ws2, [
        "ID", "Name", "Designation", "Department", "Location", "Country",
        "Status", "Availability Flag", "Availability Date",
        "Current Project", "Allocation %", "Career Level",
        "Experience (yr)", "Primary Skills", "Certifications",
    ])
    for i, e in enumerate(ctx["employees"], start=2):
        ws2.append([
            e.id, e.name, e.designation, e.department, e.location, e.country,
            e.status, _availability_flag(e, ctx["today"]), e.availability_date,
            e.current_project or "", e.allocation_pct, e.career_level,
            e.experience_years,
            ", ".join(e.primary_skills),
            ", ".join(e.certifications),
        ])
        if i % 2 == 0:
            for j in range(1, 16):
                ws2.cell(row=i, column=j).fill = _ZEBRA_FILL
    _xlsx_autosize(ws2)

    # -- Sheet 3: Projects --
    ws3 = wb.create_sheet("Projects")
    _xlsx_write_header(ws3, [
        "ID", "Name", "Client", "Status", "Priority", "Health",
        "Start Date", "End Date", "Team Size", "Required Skills",
    ])
    for i, p in enumerate(ctx["projects"], start=2):
        ws3.append([
            p.id, p.name, p.client, p.status, p.priority, p.health,
            p.start_date, p.end_date, p.team_size,
            ", ".join(p.required_skills),
        ])
        if i % 2 == 0:
            for j in range(1, 11):
                ws3.cell(row=i, column=j).fill = _ZEBRA_FILL
    _xlsx_autosize(ws3)

    # -- Sheet 4: Allocations --
    ws4 = wb.create_sheet("Allocations")
    _xlsx_write_header(ws4, [
        "ID", "Employee ID", "Project ID", "Role",
        "Allocation %", "Start Date", "End Date", "Billable",
    ])
    for i, a in enumerate(ctx["allocations"], start=2):
        ws4.append([
            a.id, a.employee_id, a.project_id, a.role,
            a.allocation_pct, a.start_date, a.end_date,
            "Yes" if a.billable else "No",
        ])
        if i % 2 == 0:
            for j in range(1, 9):
                ws4.cell(row=i, column=j).fill = _ZEBRA_FILL
    _xlsx_autosize(ws4)

    # -- Sheet 5: Skill Coverage --
    ws5 = wb.create_sheet("Skill Coverage")
    _xlsx_write_header(ws5, ["Skill", "Primary Holders"])
    for i, (skill, count) in enumerate(ctx["top_skills"], start=2):
        ws5.append([skill, count])
        if i % 2 == 0:
            for j in range(1, 3):
                ws5.cell(row=i, column=j).fill = _ZEBRA_FILL
    _xlsx_autosize(ws5)

    # -- Sheet 6: Rolling Off (30 days) --
    ws6 = wb.create_sheet("Rolling Off 30d")
    _xlsx_write_header(ws6, [
        "Name", "Designation", "Current Project", "Days Until Free", "Availability Date",
    ])
    for i, r in enumerate(ctx["rolling_off"], start=2):
        ws6.append([
            r["name"], r["designation"], r["current_project"],
            r["days_until_free"], r["availability_date"],
        ])
        if i % 2 == 0:
            for j in range(1, 6):
                ws6.cell(row=i, column=j).fill = _ZEBRA_FILL
    _xlsx_autosize(ws6)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# DOCX
# ═══════════════════════════════════════════════════════════════════════════

def _docx_heading(doc, text: str, level: int = 1) -> None:
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = DocxRGBColor(0x1F, 0x38, 0x64)


def _docx_table(doc, headers: list[str], rows: list[list[Any]]) -> None:
    tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
    tbl.style = "Light Grid Accent 1"
    for j, h in enumerate(headers):
        cell = tbl.rows[0].cells[j]
        cell.text = h
        for p in cell.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = DocxPt(10)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.rows[i].cells[j]
            cell.text = str(val)
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.size = DocxPt(9)


def generate_talent_docx_bytes(employees, projects, allocations) -> bytes:
    ctx = _build_context(employees, projects, allocations)
    doc = Document()

    # Title
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = title.add_run("Talent Management Report")
    tr.font.size = DocxPt(28)
    tr.font.bold = True
    tr.font.color.rgb = DocxRGBColor(0x1F, 0x38, 0x64)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sub.add_run(f"Generated on {ctx['today'].isoformat()}")
    sr.font.italic = True
    sr.font.size = DocxPt(11)

    doc.add_paragraph()

    # -- Team Overview --
    _docx_heading(doc, "Team Overview", level=1)
    _docx_table(doc, ["Metric", "Value"], [
        ["Total Employees", ctx["total_employees"]],
        ["Available", ctx["available"]],
        ["Allocated", ctx["allocated"]],
        ["Bench", ctx["bench"]],
        ["On Leave", ctx["on_leave"]],
        ["Rolling Off in 30 days", len(ctx["rolling_off"])],
        ["Active Projects", len(ctx["active_projects"])],
    ])

    # -- Availability breakdown --
    _docx_heading(doc, "Availability Breakdown", level=1)
    rows = [
        [e.name, e.designation, e.status.upper(), _availability_flag(e, ctx["today"])]
        for e in ctx["employees"][:20]
    ]
    _docx_table(doc, ["Name", "Designation", "Status", "Availability"], rows)

    # -- Skill Coverage --
    _docx_heading(doc, "Skill Coverage (top 15)", level=1)
    _docx_table(doc, ["Skill", "Primary Holders"], [[s, c] for s, c in ctx["top_skills"]])

    # -- Active Projects --
    _docx_heading(doc, "Active Projects", level=1)
    proj_rows = [
        [p.name, p.client, p.status.upper(), p.end_date]
        for p in ctx["active_projects"]
    ]
    _docx_table(doc, ["Project", "Client", "Status", "Ends"],
                proj_rows or [["No active projects", "-", "-", "-"]])

    # -- Rolling off --
    _docx_heading(doc, "Rolling Off in Next 30 Days", level=1)
    ro_rows = [
        [r["name"], r["designation"], r["current_project"],
         f"{r['days_until_free']}d ({r['availability_date']})"]
        for r in ctx["rolling_off"]
    ]
    _docx_table(doc, ["Name", "Designation", "Current Project", "Free In"],
                ro_rows or [["None", "-", "-", "-"]])

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
