"""
PPT Generator — builds a professional PM status presentation from Jira data.

Uses python-pptx (pure Python, no Office required).

Slide deck structure:
  1. Title slide
  2. Executive Summary  (portfolio table — all projects at a glance)
  3–N. Per-project pair: Status slide + Risk slide
  N+1. Next Steps / Recommendations

Call:
    path = generate_ppt(projects_data, risks_data, output_path)

where:
    projects_data — list of dicts from jira_client.get_project_status()
    risks_data    — list of dicts from jira_client.get_risks()
    output_path   — file path to write the .pptx to
"""
from __future__ import annotations

import io
import os
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
C_DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)   # slide headers / title bg
C_MED_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # accent / metric boxes
C_LIGHT_BG   = RGBColor(0xD6, 0xE4, 0xF0)   # alternating table rows
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN      = RGBColor(0x70, 0xAD, 0x47)   # done / healthy
C_AMBER      = RGBColor(0xFF, 0xC0, 0x00)   # at-risk / medium
C_RED        = RGBColor(0xC0, 0x00, 0x00)   # overdue / high-risk
C_DARK_TEXT  = RGBColor(0x26, 0x26, 0x26)
C_GREY_TEXT  = RGBColor(0x60, 0x60, 0x60)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

HEADER_H   = Inches(1.0)
MARGIN_L   = Inches(0.5)
MARGIN_T   = Inches(1.1)
CONTENT_W  = SLIDE_W - MARGIN_L * 2


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _bg(slide_or_shape, color: RGBColor) -> None:
    fill = slide_or_shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(
    slide,
    left, top, width, height,
    text: str,
    font_size: int = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = C_DARK_TEXT,
    bg: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        _bg(txBox, bg)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _slide_header(slide, title: str) -> None:
    """Dark-blue banner at the top of a content slide."""
    banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, HEADER_H,
    )
    _bg(banner, C_DARK_BLUE)
    banner.line.fill.background()  # no border
    _textbox(
        slide,
        MARGIN_L, Inches(0.15), CONTENT_W, Inches(0.7),
        title,
        font_size=22, bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )


def _cell_bg(cell, color: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _cell_text(cell, text: str, font_size: int = 10, bold: bool = False,
               color: RGBColor = C_DARK_TEXT, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else "—"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _metric_box(slide, left, top, width, height,
                label: str, value: str, value_color: RGBColor = C_DARK_BLUE) -> None:
    """Coloured KPI tile."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    _bg(box, C_LIGHT_BG)
    box.line.color.rgb = C_MED_BLUE
    # value
    _textbox(slide, left + Inches(0.1), top + Inches(0.05),
             width - Inches(0.2), Inches(0.45),
             value, font_size=24, bold=True, color=value_color,
             align=PP_ALIGN.CENTER)
    # label
    _textbox(slide, left + Inches(0.05), top + Inches(0.5),
             width - Inches(0.1), Inches(0.3),
             label, font_size=9, color=C_GREY_TEXT,
             align=PP_ALIGN.CENTER)


def _severity_color(severity: str) -> RGBColor:
    return {
        "HIGH": C_RED,
        "MEDIUM": C_AMBER,
        "LOW": C_GREEN,
    }.get(severity.upper(), C_DARK_TEXT)


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # truly blank
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _blank_slide(prs)

    # Full-slide dark-blue background
    bg_rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    _bg(bg_rect, C_DARK_BLUE)
    bg_rect.line.fill.background()

    # Accent bar (medium blue strip)
    bar = slide.shapes.add_shape(1, 0, Inches(4.0), SLIDE_W, Inches(0.08))
    _bg(bar, C_MED_BLUE)
    bar.line.fill.background()

    # Main title
    _textbox(slide, Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
             title, font_size=40, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    # Subtitle / date
    _textbox(slide, Inches(1), Inches(3.2), Inches(11.33), Inches(0.7),
             subtitle, font_size=18, color=C_LIGHT_BG,
             align=PP_ALIGN.CENTER)

    # Footer tag
    _textbox(slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
             "Generated by PM Delivery Agent  •  Confidential",
             font_size=10, italic=True, color=C_GREY_TEXT,
             align=PP_ALIGN.CENTER)


def _section_divider(prs: Presentation, section_title: str) -> None:
    slide = _blank_slide(prs)

    left_bar = slide.shapes.add_shape(1, 0, 0, Inches(0.25), SLIDE_H)
    _bg(left_bar, C_MED_BLUE)
    left_bar.line.fill.background()

    _textbox(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(1.2),
             section_title, font_size=32, bold=True, color=C_DARK_BLUE,
             align=PP_ALIGN.LEFT)


def _exec_summary_slide(prs: Presentation, all_status: list[dict], all_risks: list[dict]) -> None:
    slide = _blank_slide(prs)
    _slide_header(slide, "Executive Summary — Portfolio Overview")

    # Build risk lookup
    risk_map = {r.get("project_key", ""): r for r in all_risks}

    # Table: project | total | done | in progress | overdue | risk
    rows = 1 + len(all_status)
    cols = 6
    top = MARGIN_T
    tbl = slide.shapes.add_table(rows, cols, MARGIN_L, top, CONTENT_W, Inches(0.38 * rows)).table

    headers = ["Project", "Total Issues", "Done", "In Progress", "Overdue", "Risk Level"]
    col_widths = [Inches(3.0), Inches(1.5), Inches(1.3), Inches(1.7), Inches(1.4), Inches(1.4)]
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        tbl.columns[i].width = w
        _cell_bg(tbl.cell(0, i), C_DARK_BLUE)
        _cell_text(tbl.cell(0, i), h, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for row_i, status in enumerate(all_status, 1):
        pk = status.get("project_key", "")
        risk = risk_map.get(pk, {})
        severity = risk.get("severity", "—")
        overdue_count = len(status.get("overdue", []))
        bg = C_WHITE if row_i % 2 == 0 else C_LIGHT_BG
        row_data = [
            f"{status.get('project_name', pk)} ({pk})",
            str(status.get("total", 0)),
            str(status.get("done", 0)),
            str(status.get("in_progress", 0)),
            str(overdue_count),
            severity,
        ]
        row_colors = [bg, bg, bg, bg,
                      C_RED if overdue_count > 0 else bg,
                      _severity_color(severity) if severity != "—" else bg]
        text_colors = [C_DARK_TEXT, C_DARK_TEXT, C_DARK_TEXT, C_DARK_TEXT,
                       C_WHITE if overdue_count > 0 else C_DARK_TEXT,
                       C_WHITE if severity in ("HIGH", "MEDIUM") else C_DARK_TEXT]
        for col_i, (val, bc, tc) in enumerate(zip(row_data, row_colors, text_colors)):
            _cell_bg(tbl.cell(row_i, col_i), bc)
            _cell_text(tbl.cell(row_i, col_i), val, font_size=10, color=tc,
                       align=PP_ALIGN.CENTER if col_i > 0 else PP_ALIGN.LEFT)


def _project_status_slide(prs: Presentation, status: dict) -> None:
    pk = status.get("project_key", "")
    name = status.get("project_name", pk)
    slide = _blank_slide(prs)
    _slide_header(slide, f"Project Status — {name}  [{pk}]")

    total    = status.get("total", 0)
    done     = status.get("done", 0)
    in_prog  = status.get("in_progress", 0)
    in_rev   = status.get("in_review", 0)
    todo     = status.get("todo", 0)
    overdue  = len(status.get("overdue", []))
    pct_done = int(done / total * 100) if total else 0

    # KPI metric boxes
    box_w = Inches(1.8)
    box_h = Inches(0.85)
    box_top = MARGIN_T
    metrics = [
        ("Total Issues",  str(total),          C_MED_BLUE),
        ("Done",          f"{done} ({pct_done}%)", C_GREEN),
        ("In Progress",   str(in_prog),         C_MED_BLUE),
        ("In Review",     str(in_rev),          C_MED_BLUE),
        ("To Do",         str(todo),            C_GREY_TEXT),
        ("Overdue",       str(overdue),         C_RED if overdue else C_GREEN),
    ]
    for i, (label, val, col) in enumerate(metrics):
        _metric_box(slide,
                    MARGIN_L + i * (box_w + Inches(0.12)), box_top,
                    box_w, box_h, label, val, col)

    # Progress bar
    bar_top = box_top + box_h + Inches(0.15)
    bar_w   = CONTENT_W
    bar_h   = Inches(0.18)
    bg_bar  = slide.shapes.add_shape(1, MARGIN_L, bar_top, bar_w, bar_h)
    _bg(bg_bar, C_LIGHT_BG)
    bg_bar.line.fill.background()
    if total and done:
        fill_w = int(bar_w * done / total)
        fg_bar = slide.shapes.add_shape(1, MARGIN_L, bar_top, fill_w, bar_h)
        _bg(fg_bar, C_GREEN)
        fg_bar.line.fill.background()
    _textbox(slide, MARGIN_L, bar_top, bar_w, bar_h,
             f"  {pct_done}% complete", font_size=9, bold=True, color=C_DARK_TEXT)

    # Issue list table
    issues = status.get("issues", [])[:12]
    if not issues:
        _textbox(slide, MARGIN_L, bar_top + Inches(0.3), CONTENT_W, Inches(0.4),
                 "No issues found in this project.", font_size=11, color=C_GREY_TEXT)
        return

    tbl_top = bar_top + bar_h + Inches(0.18)
    rows = 1 + len(issues)
    tbl = slide.shapes.add_table(rows, 5, MARGIN_L, tbl_top, CONTENT_W,
                                 Inches(0.32 * rows)).table
    headers = ["Key", "Summary", "Status", "Priority", "Due Date"]
    col_widths = [Inches(1.3), Inches(5.5), Inches(1.6), Inches(1.3), Inches(1.5)]
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        tbl.columns[i].width = w
        _cell_bg(tbl.cell(0, i), C_MED_BLUE)
        _cell_text(tbl.cell(0, i), h, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    today = date.today().isoformat()
    for row_i, issue in enumerate(issues, 1):
        is_overdue = bool(issue.get("due_date") and issue.get("due_date") < today
                          and (issue.get("status") or "").lower() != "done")
        bg = C_RED if is_overdue else (C_WHITE if row_i % 2 == 0 else C_LIGHT_BG)
        tc = C_WHITE if is_overdue else C_DARK_TEXT
        row_data = [
            issue.get("key", ""),
            issue.get("summary", ""),
            issue.get("status", ""),
            issue.get("priority", ""),
            issue.get("due_date", "—") or "—",
        ]
        for col_i, val in enumerate(row_data):
            _cell_bg(tbl.cell(row_i, col_i), bg)
            _cell_text(tbl.cell(row_i, col_i), val, font_size=9, color=tc)

    if len(status.get("issues", [])) > 12:
        _textbox(slide, MARGIN_L, tbl_top + Inches(0.32 * rows) + Inches(0.05),
                 CONTENT_W, Inches(0.25),
                 f"  + {len(status['issues']) - 12} more issues not shown",
                 font_size=8, italic=True, color=C_GREY_TEXT)


def _project_risk_slide(prs: Presentation, risk: dict) -> None:
    pk   = risk.get("project_key", "")
    name = risk.get("project_name", pk)
    slide = _blank_slide(prs)
    severity   = risk.get("severity", "LOW")
    risk_count = risk.get("risk_count", 0)
    sev_color  = _severity_color(severity)

    _slide_header(slide, f"Risk Report — {name}  [{pk}]")

    # Severity badge
    badge_w = Inches(2.5)
    badge_h = Inches(0.9)
    badge = slide.shapes.add_shape(1, MARGIN_L, MARGIN_T, badge_w, badge_h)
    _bg(badge, sev_color)
    badge.line.fill.background()
    _textbox(slide, MARGIN_L, MARGIN_T + Inches(0.05), badge_w, Inches(0.38),
             severity, font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, MARGIN_T + Inches(0.45), badge_w, Inches(0.3),
             f"{risk_count} HIGH/HIGHEST issue{'s' if risk_count != 1 else ''}",
             font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

    risks = risk.get("risks", [])
    if not risks:
        _textbox(slide, MARGIN_L, MARGIN_T + Inches(1.1), CONTENT_W, Inches(0.5),
                 "No high-priority risks found. Project health looks good.",
                 font_size=13, color=C_GREEN, bold=True)
        return

    rows = 1 + len(risks)
    tbl_top = MARGIN_T + badge_h + Inches(0.2)
    tbl = slide.shapes.add_table(rows, 5, MARGIN_L, tbl_top, CONTENT_W,
                                 Inches(0.35 * rows)).table
    headers = ["Key", "Summary", "Priority", "Due Date", "Overdue?"]
    col_widths = [Inches(1.3), Inches(5.8), Inches(1.4), Inches(1.5), Inches(1.3)]
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        tbl.columns[i].width = w
        _cell_bg(tbl.cell(0, i), C_DARK_BLUE)
        _cell_text(tbl.cell(0, i), h, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for row_i, r in enumerate(risks, 1):
        is_overdue = r.get("overdue", False)
        bg = C_RED if is_overdue else (C_WHITE if row_i % 2 == 0 else C_LIGHT_BG)
        tc = C_WHITE if is_overdue else C_DARK_TEXT
        row_data = [
            r.get("key", ""),
            r.get("summary", ""),
            r.get("priority", ""),
            r.get("due_date", "—") or "—",
            "YES  ⚠" if is_overdue else "No",
        ]
        for col_i, val in enumerate(row_data):
            _cell_bg(tbl.cell(row_i, col_i), bg)
            _cell_text(tbl.cell(row_i, col_i), val, font_size=9, color=tc,
                       align=PP_ALIGN.CENTER if col_i > 1 else PP_ALIGN.LEFT)


def _next_steps_slide(prs: Presentation, all_status: list[dict], all_risks: list[dict]) -> None:
    slide = _blank_slide(prs)
    _slide_header(slide, "Recommended Next Steps")

    steps: list[str] = []
    for risk in all_risks:
        pk = risk.get("project_key", "")
        name = risk.get("project_name", pk)
        if risk.get("severity") == "HIGH":
            steps.append(f"[{pk}] {name} — Immediately address {risk.get('risk_count')} high-priority blockers.")
    for status in all_status:
        pk   = status.get("project_key", "")
        name = status.get("project_name", pk)
        overdue = status.get("overdue", [])
        if overdue:
            steps.append(f"[{pk}] {name} — Resolve {len(overdue)} overdue issue(s); review deadlines with team.")
    for status in all_status:
        pk   = status.get("project_key", "")
        name = status.get("project_name", pk)
        total = status.get("total", 0)
        done  = status.get("done", 0)
        if total and done / total < 0.5:
            steps.append(f"[{pk}] {name} — Less than 50% complete; consider resourcing review.")

    if not steps:
        steps.append("All projects appear healthy. Continue regular sprint reviews.")
        steps.append("Schedule stakeholder sync to validate upcoming milestones.")
        steps.append("Ensure all issues have assigned owners and due dates.")

    steps.append("Share this report with stakeholders via the PM Agent Slack channel.")

    top = MARGIN_T
    for i, step in enumerate(steps[:8]):
        bullet = slide.shapes.add_shape(1,
                                        MARGIN_L, top + i * Inches(0.55),
                                        Inches(0.12), Inches(0.35))
        _bg(bullet, C_MED_BLUE)
        bullet.line.fill.background()
        _textbox(slide, MARGIN_L + Inches(0.25), top + i * Inches(0.55),
                 CONTENT_W - Inches(0.3), Inches(0.45),
                 step, font_size=12, color=C_DARK_TEXT)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def generate_ppt(
    projects_data: list[dict[str, Any]],
    risks_data: list[dict[str, Any]],
    output_path: str,
) -> str:
    """
    Build a .pptx file from Jira data and write it to output_path.
    Returns output_path on success.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    today = date.today().strftime("%B %d, %Y")
    project_names = ", ".join(
        p.get("project_name", p.get("project_key", "")) for p in projects_data
    ) or "All Projects"

    # 1. Title slide
    _title_slide(
        prs,
        title=f"PM Portfolio Report",
        subtitle=f"{project_names}  •  {today}",
    )

    # 2. Executive Summary
    _exec_summary_slide(prs, projects_data, risks_data)

    # Build risk lookup
    risk_map = {r.get("project_key", ""): r for r in risks_data}

    # 3–N: Per-project slides
    for status in projects_data:
        pk = status.get("project_key", "")
        _section_divider(prs, f"{status.get('project_name', pk)}  [{pk}]")
        _project_status_slide(prs, status)
        risk = risk_map.get(pk, {"project_key": pk, "project_name": status.get("project_name", pk),
                                  "severity": "LOW", "risk_count": 0, "risks": []})
        _project_risk_slide(prs, risk)

    # Last slide: Next Steps
    _next_steps_slide(prs, projects_data, risks_data)

    prs.save(output_path)
    return output_path


def generate_ppt_bytes(
    projects_data: list[dict[str, Any]],
    risks_data: list[dict[str, Any]],
) -> bytes:
    """Generate PPT and return raw bytes (useful for streaming HTTP response)."""
    buf = io.BytesIO()
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    today = date.today().strftime("%B %d, %Y")
    project_names = ", ".join(
        p.get("project_name", p.get("project_key", "")) for p in projects_data
    ) or "All Projects"

    _title_slide(prs, "PM Portfolio Report", f"{project_names}  •  {today}")
    _exec_summary_slide(prs, projects_data, risks_data)
    risk_map = {r.get("project_key", ""): r for r in risks_data}
    for status in projects_data:
        pk = status.get("project_key", "")
        _section_divider(prs, f"{status.get('project_name', pk)}  [{pk}]")
        _project_status_slide(prs, status)
        risk = risk_map.get(pk, {"project_key": pk, "project_name": status.get("project_name", pk),
                                  "severity": "LOW", "risk_count": 0, "risks": []})
        _project_risk_slide(prs, risk)
    _next_steps_slide(prs, projects_data, risks_data)

    prs.save(buf)
    return buf.getvalue()
